from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import os



def extract_pdf(file_path):
    reader = PdfReader(file_path)

    text = ""

    for page in reader.pages:
        text += page.extract_text() or ""

    return text


def extract_docx(file_path):
    document = Document(file_path)

    text = ""

    for paragraph in document.paragraphs:
        text += paragraph.text + "\n"

    return text


def extract_pptx(file_path):
    presentation = Presentation(file_path)

    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text



def clean_text(text):
    text = " ".join(text.split())
    return text.strip()




def categorize_document(file_path, text=""):
    content = f"{file_path} {text}".lower()

    if (
        "question bank" in content
        or "question paper" in content
        or "question" in content
        or "qp" in content
        or "exam paper" in content
    ):
        return "Question Paper"

    elif "lab manual" in content or "lab" in content or "practical" in content:
        return "Lab Manual"

    elif "textbook" in content or "text book" in content:
        return "Textbook"

    elif "assignment" in content:
        return "Assignment"

    elif "lecture" in content or "notes" in content or "note" in content:
        return "Notes"

    else:
        return "General"


# -----------------------------
# File Type
# -----------------------------

def detect_document_type(file_path):
    filename = file_path.lower()

    if filename.endswith(".pdf"):
        return "PDF"

    elif filename.endswith(".docx"):
        return "DOCX"

    elif filename.endswith(".pptx"):
        return "PPTX"

    else:
        return "Unsupported"




SUBJECT_KEYWORDS = {
    "DBMS": [
        "dbms",
        "database management",
        "database",
        "sql",
        "mysql",
        "postgresql",
        "oracle",
        "normalization",
        "relational database",
        "primary key",
        "foreign key",
        "er diagram",
        "transaction",
    ],

    "Data Structures and Algorithms": [
        "data structure",
        "algorithm",
        "array",
        "linked list",
        "stack",
        "queue",
        "tree",
        "binary tree",
        "graph",
        "sorting",
        "searching",
        "hashing",
        "recursion",
    ],

    "Operating Systems": [
        "operating system",
        "process",
        "thread",
        "deadlock",
        "cpu scheduling",
        "process scheduling",
        "memory management",
        "paging",
        "segmentation",
        "virtual memory",
        "file system",
    ],

    "Computer Networks": [
        "computer network",
        "networking",
        "tcp",
        "udp",
        "ip address",
        "routing",
        "osi model",
        "tcp/ip",
        "subnetting",
        "network protocol",
        "ethernet",
    ],

    "Aptitude": [
        "aptitude",
        "quantitative aptitude",
        "logical reasoning",
        "verbal ability",
        "engineering aptitude",
        "comprehensive engineering aptitude",
        "percentage",
        "profit and loss",
        "ratio",
        "proportion",
        "time and work",
        "probability",
        "permutation",
        "combination",
        "number system",
    ],

    "Artificial Intelligence and Machine Learning": [
        "artificial intelligence",
        "machine learning",
        "deep learning",
        "neural network",
        "cnn",
        "rnn",
        "classification",
        "regression",
        "supervised learning",
        "unsupervised learning",
        "computer vision",
    ],
}


def classify_subject(file_path, text=""):
    content = f"{file_path} {text}".lower()

    subject_scores = {}

    for subject, keywords in SUBJECT_KEYWORDS.items():
        score = 0

        for keyword in keywords:
            if keyword in content:
                score += 1

        if score > 0:
            subject_scores[subject] = score

    if not subject_scores:
        return "General"

    return max(subject_scores, key=subject_scores.get)



TOPIC_KEYWORDS = {
    "DBMS": {
        "SQL": ["sql", "select", "insert", "update", "delete", "query"],
        "Normalization": [
            "normalization",
            "1nf",
            "2nf",
            "3nf",
            "bcnf",
        ],
        "ER Model": [
            "er diagram",
            "entity relationship",
            "entity",
            "relationship",
        ],
        "Database Keys": [
            "primary key",
            "foreign key",
            "candidate key",
            "super key",
        ],
        "Transactions": [
            "transaction",
            "acid",
            "commit",
            "rollback",
            "concurrency",
        ],
    },

    "Data Structures and Algorithms": {
        "Arrays": ["array"],
        "Linked Lists": ["linked list"],
        "Stacks and Queues": ["stack", "queue"],
        "Trees": ["tree", "binary tree", "bst"],
        "Graphs": ["graph", "bfs", "dfs"],
        "Sorting": ["sorting", "bubble sort", "merge sort", "quick sort"],
        "Searching": ["searching", "binary search", "linear search"],
    },

    "Operating Systems": {
        "Process Management": ["process", "thread"],
        "CPU Scheduling": [
            "cpu scheduling",
            "fcfs",
            "sjf",
            "round robin",
            "priority scheduling",
        ],
        "Deadlocks": ["deadlock", "banker's algorithm"],
        "Memory Management": [
            "memory management",
            "paging",
            "segmentation",
            "virtual memory",
        ],
        "File Systems": ["file system", "file allocation"],
    },

    "Computer Networks": {
        "OSI Model": ["osi model", "osi"],
        "TCP/IP": ["tcp", "udp", "tcp/ip"],
        "IP Addressing": ["ip address", "ipv4", "ipv6", "subnetting"],
        "Routing": ["routing", "router", "routing protocol"],
        "Network Protocols": ["http", "https", "ftp", "dns", "dhcp"],
    },

    "Aptitude": {
        "Percentages": ["percentage", "percent"],
        "Ratio and Proportion": ["ratio", "proportion"],
        "Profit and Loss": ["profit", "loss", "profit and loss"],
        "Time and Work": ["time and work", "work and time"],
        "Probability": ["probability"],
        "Number System": ["number system", "divisibility"],
        "Logical Reasoning": [
            "logical reasoning",
            "reasoning",
            "coding decoding",
            "blood relation",
            "syllogism",
        ],
        "Verbal Ability": [
            "verbal ability",
            "grammar",
            "synonym",
            "antonym",
            "sentence correction",
        ],
    },
}


def classify_topic(subject, file_path, text=""):
    content = f"{file_path} {text}".lower()

    topics = TOPIC_KEYWORDS.get(subject, {})

    topic_scores = {}

    for topic, keywords in topics.items():
        score = 0

        for keyword in keywords:
            if keyword in content:
                score += 1

        if score > 0:
            topic_scores[topic] = score

    if not topic_scores:
        return "General"

    return max(topic_scores, key=topic_scores.get)



def process_document(file_path, original_filename=None):

    try:
        if file_path.lower().endswith(".pdf"):
            text = extract_pdf(file_path)

        elif file_path.lower().endswith(".docx"):
            text = extract_docx(file_path)

        elif file_path.lower().endswith(".pptx"):
            text = extract_pptx(file_path)

        else:
            return {
                "error": "Unsupported file type"
            }

        text = clean_text(text)

        classification_name = original_filename or os.path.basename(file_path)

        if not text:
            return {
                "filename": classification_name,
                "file_type": detect_document_type(file_path),
                "category": "Unknown",
                "subject": "Unknown",
                "topic": "Unknown",
                "text": "",
                "error": "The document is empty or contains no readable text"
            }

        category = categorize_document(classification_name, text)
        subject = classify_subject(classification_name, text)
        topic = classify_topic(subject, classification_name, text)

        return {
            "filename": classification_name,
            "file_type": detect_document_type(file_path),
            "category": category,
            "subject": subject,
            "topic": topic,
            "text": text
        }

    except Exception as e:
        return {
            "filename": original_filename or os.path.basename(file_path),
            "file_type": detect_document_type(file_path),
            "category": "Unknown",
            "subject": "Unknown",
            "topic": "Unknown",
            "text": "",
            "error": "The document could not be processed"
        }

if __name__ == "__main__":

    result = process_document(
        "../samples/Question Bank for Comprehensive Engineering Aptitude Test.pdf"
    )

    print("Filename:", result["filename"])
    print("File Type:", result["file_type"])
    print("Category:", result["category"])
    print("Subject:", result["subject"])
    print("Topic:", result["topic"])
    print("Text:")
    print(result["text"])