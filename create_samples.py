from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from docx import Document
from pptx import Presentation


# ---------------- PDF ----------------

pdf_file = "samples/DBMS_Notes.pdf"

styles = getSampleStyleSheet()
doc = SimpleDocTemplate(pdf_file, pagesize=A4)

content = [
    Paragraph("Database Management Systems - Lecture Notes", styles["Title"]),
    Spacer(1, 20),

    Paragraph("<b>1. Introduction to DBMS</b>", styles["Heading2"]),
    Paragraph(
        "A Database Management System (DBMS) is software used to create, "
        "store, organize, retrieve and manage data efficiently. Examples "
        "include MySQL, PostgreSQL, Oracle and Microsoft SQL Server.",
        styles["BodyText"]
    ),
    Spacer(1, 10),

    Paragraph("<b>2. Database Models</b>", styles["Heading2"]),
    Paragraph(
        "Common database models include hierarchical, network, relational "
        "and object-oriented models. The relational model stores data in "
        "tables consisting of rows and columns.",
        styles["BodyText"]
    ),
    Spacer(1, 10),

    Paragraph("<b>3. Normalization</b>", styles["Heading2"]),
    Paragraph(
        "Normalization is the process of organizing data to reduce "
        "redundancy and improve data integrity. First Normal Form (1NF) "
        "requires atomic values. Second Normal Form (2NF) removes partial "
        "dependencies. Third Normal Form (3NF) removes transitive "
        "dependencies.",
        styles["BodyText"]
    ),
    Spacer(1, 10),

    Paragraph("<b>4. Primary and Foreign Keys</b>", styles["Heading2"]),
    Paragraph(
        "A primary key uniquely identifies each record in a table. "
        "A foreign key creates a relationship between two tables by "
        "referencing a primary key in another table.",
        styles["BodyText"]
    ),
    Spacer(1, 10),

    Paragraph("<b>5. SQL</b>", styles["Heading2"]),
    Paragraph(
        "SQL is used to interact with relational databases. Important "
        "commands include SELECT, INSERT, UPDATE and DELETE. The SELECT "
        "statement is commonly used to retrieve records from a table.",
        styles["BodyText"]
    ),
]

doc.build(content)


# ---------------- DOCX ----------------

docx_file = "samples/DBMS_Lab_Notes.docx"

document = Document()

document.add_heading("Database Management Systems - Lab Notes", 0)

document.add_heading("Experiment 1: Creating a Database", level=1)
document.add_paragraph(
    "Objective: To create a database and understand the basic structure "
    "of tables using SQL."
)

document.add_paragraph(
    "SQL Command:\n"
    "CREATE DATABASE CollegeDB;"
)

document.add_heading("Experiment 2: Creating a Student Table", level=1)
document.add_paragraph(
    "The following table can be used to store student information."
)

document.add_paragraph(
    "CREATE TABLE Student (\n"
    "    Student_ID INT PRIMARY KEY,\n"
    "    Student_Name VARCHAR(50),\n"
    "    Department VARCHAR(30),\n"
    "    CGPA DECIMAL(3,2)\n"
    ");"
)

document.add_heading("Experiment 3: Retrieving Records", level=1)
document.add_paragraph(
    "The SELECT statement retrieves records from a database table."
)

document.add_paragraph(
    "Example:\n"
    "SELECT * FROM Student;"
)

document.add_heading("Result", level=1)
document.add_paragraph(
    "The database and Student table are created successfully, and records "
    "can be retrieved using SQL queries."
)

document.save(docx_file)


# ---------------- PPTX ----------------

pptx_file = "samples/DBMS_Lecture.pptx"

presentation = Presentation()

slides = [
    (
        "Database Management Systems",
        "Introduction to DBMS\nDatabase models\nRelational databases\nSQL basics"
    ),
    (
        "What is a DBMS?",
        "A DBMS is software used to manage databases.\n"
        "It provides data storage and retrieval.\n"
        "It improves security and data consistency.\n"
        "Examples: MySQL, PostgreSQL and Oracle."
    ),
    (
        "Database Normalization",
        "Normalization reduces data redundancy.\n"
        "1NF removes repeating groups.\n"
        "2NF removes partial dependency.\n"
        "3NF removes transitive dependency."
    ),
    (
        "Keys in a Database",
        "Primary key uniquely identifies a record.\n"
        "Foreign key connects related tables.\n"
        "Candidate keys can uniquely identify records."
    ),
    (
        "SQL Commands",
        "SELECT - retrieve data\n"
        "INSERT - add data\n"
        "UPDATE - modify data\n"
        "DELETE - remove data"
    ),
]

for title, body in slides:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[1]
    )
    slide.shapes.title.text = title
    slide.placeholders[1].text = body

presentation.save(pptx_file)

print("Sample academic files created successfully.")
